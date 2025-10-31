"""File parsing utilities for uploaded documents."""

from __future__ import annotations

import io
import re
from pathlib import Path
from typing import Iterable

from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation
from docx import Document


class UnsupportedFileTypeError(ValueError):
    """Raised when the uploaded file type is not supported."""


class FileParser:
    """Parse supported document formats to plain text."""

    SUPPORTED_SUFFIXES = {".pdf", ".docx", ".pptx", ".txt", ".md"}

    def parse_bytes(self, filename: str, content: bytes) -> str:
        suffix = Path(filename).suffix.lower()
        if suffix not in self.SUPPORTED_SUFFIXES:
            raise UnsupportedFileTypeError(f"Unsupported file type: {suffix}")

        if suffix == ".pdf":
            text = self._parse_pdf(content)
        elif suffix == ".docx":
            text = self._parse_docx(content)
        elif suffix == ".pptx":
            text = self._parse_pptx(content)
        else:  # .txt or .md
            text = self._parse_plaintext(content)
        return self._normalize_text(text)

    def parse_batch(self, files: Iterable[tuple[str, bytes]]) -> list[str]:
        return [self.parse_bytes(name, data) for name, data in files]

    def _parse_pdf(self, content: bytes) -> str:
        buffer = io.BytesIO(content)
        try:
            return pdf_extract_text(buffer) or ""
        except Exception:  # pragma: no cover - pdfminer raises many error types
            return ""

    def _parse_docx(self, content: bytes) -> str:
        buffer = io.BytesIO(content)
        try:
            document = Document(buffer)
        except Exception:  # pragma: no cover - document parsing errors
            return ""
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        return "\n".join(paragraphs)

    def _parse_pptx(self, content: bytes) -> str:
        buffer = io.BytesIO(content)
        try:
            presentation = Presentation(buffer)
        except Exception:  # pragma: no cover
            return ""
        texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
        return "\n".join(texts)

    def _parse_plaintext(self, content: bytes) -> str:
        for encoding in ("utf-8", "latin-1"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return ""

    def _normalize_text(self, text: str) -> str:
        cleaned = re.sub(r"\s+", " ", text).strip()
        return cleaned
